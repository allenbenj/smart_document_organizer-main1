"""
Content Extraction Module

Extracts text and metadata from various file types for better organization.
"""

import re
import logging
from pathlib import Path
from typing import Dict, List, Optional, Tuple, Any
from datetime import datetime

from file_organizer.services.llm_entity_extractor import ExtractedEntity, LLMEntityExtractor
from file_organizer.services.service_container import ServiceContainer

logger = logging.getLogger(__name__)


class ContentExtractor:
    """
    Extract content and metadata from files for intelligent organization.
    
    Supports:
    - PDF (using pdfplumber or PyPDF2)
    - DOCX (using python-docx)
    - Plain text (.txt, .md, .log)
    - HTML (basic text extraction)
    """
    
    # Document type keywords to look for in content
    DOC_TYPE_KEYWORDS = {
        'motion': ['motion to', 'moves the court', 'motion for', 'hereby moves'],
        'order': ['it is ordered', 'court orders', 'hereby ordered', 'order granting', 'order denying'],
        'complaint': ['complaint', 'plaintiff alleges', 'cause of action', 'jurisdiction'],
        'affidavit': ['affidavit', 'sworn statement', 'under penalty of perjury', 'affiant states'],
        'subpoena': ['subpoena', 'commanded to appear', 'produce documents', 'witness'],
        'brief': ['brief in support', 'memorandum of law', 'argument', 'conclusion'],
        'petition': ['petition for', 'petitioner', 'prays the court', 'petitioner requests'],
        'notice': ['notice of', 'hereby notified', 'take notice', 'notice is given'],
        'discovery': ['interrogatories', 'request for production', 'request for admission', 'deposition notice'],
        'exhibit': ['exhibit', 'attached hereto', 'marked as exhibit'],
        'deposition': ['deposition of', 'deponent', 'q:', 'a:', 'examination'],
        'warrant': ['warrant', 'probable cause', 'search warrant', 'arrest warrant'],
        'indictment': ['indictment', 'grand jury', 'true bill', 'count'],
        'judgment': ['judgment', 'judgment is entered', 'judgment for', 'final judgment'],
        'settlement': ['settlement', 'agree to settle', 'settlement agreement', 'release'],
        'contract': ['agreement', 'contract', 'parties agree', 'terms and conditions'],
        'letter': ['dear', 'sincerely', 'regards', 'please find enclosed'],
        'report': ['report', 'analysis', 'findings', 'conclusion', 'summary'],
        'invoice': ['invoice', 'amount due', 'payment', 'bill to', 'total'],
        'receipt': ['receipt', 'received from', 'payment received', 'thank you'],
    }
    
    # Common legal entities to extract
    ENTITY_PATTERNS = {
        'case_number': [
            r'(?:Case|No\.|#|Cause)\s*(?:No\.?)?\s*:?\s*(\d{2,4}[-]?[A-Z]{0,3}[-]?\d{3,7})',
            r'\b(\d{2}[A-Z]{2}\d{5,7})\b',
            r'(?:Docket|File)\s*(?:No\.?)?\s*:?\s*([A-Z0-9-]+)',
        ],
        'court': [
            r'(?:IN THE\s+)?([A-Z][A-Z\s]+(?:COURT|TRIBUNAL))',
            r'((?:Superior|District|Circuit|Family|Federal|Supreme|Municipal)\s+Court)',
        ],
        'judge': [
            r'(?:Judge|Hon\.|Honorable)\s+([A-Z][a-z]+\s+[A-Z][a-z]+)',
            r'(?:JUDGE|HON\.)\s+([A-Z\s]+)',
        ],
        'attorney': [
            r'(?:Attorney|Counsel|Esq\.?)\s*:?\s*([A-Z][a-z]+\s+[A-Z][a-z]+)',
            r'([A-Z][a-z]+\s+[A-Z][a-z]+),?\s+(?:Esq\.|Attorney|Counsel)',
        ],
        'date': [
            r'\b(\d{1,2}/\d{1,2}/\d{2,4})\b',
            r'\b(\d{4}-\d{2}-\d{2})\b',
            r'\b((?:January|February|March|April|May|June|July|August|September|October|November|December)\s+\d{1,2},?\s+\d{4})\b',
        ],
        'parties': [
            r'([A-Z][A-Z\s]+)\s+(?:v\.|vs\.?|versus)\s+([A-Z][A-Z\s]+)',
            r'(?:Plaintiff|Petitioner)[:\s]+([A-Z][a-zA-Z\s,]+)',
            r'(?:Defendant|Respondent)[:\s]+([A-Z][a-zA-Z\s,]+)',
        ],
        'amount': [
            r'\$\s*([\d,]+\.?\d*)',
            r'(\d{1,3}(?:,\d{3})+(?:\.\d{2})?)\s*(?:dollars|USD)',
        ],
        'phone': [
            r'\b(\d{3}[-.\s]?\d{3}[-.\s]?\d{4})\b',
            r'\((\d{3})\)\s*(\d{3})[-.\s]?(\d{4})',
        ],
        'email': [
            r'\b([a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,})\b',
        ],
    }
    
    def __init__(
        self,
        service_container: ServiceContainer | None = None,
        llm_entity_extractor: LLMEntityExtractor | None = None,
    ):
        self._pdf_extractor = None
        self._docx_available = False
        self._service_container = service_container
        if llm_entity_extractor is not None:
            self._llm_entity_extractor = llm_entity_extractor
        elif service_container is not None:
            self._llm_entity_extractor = LLMEntityExtractor(service_container)
        else:
            self._llm_entity_extractor = None
        self._init_extractors()

    def set_service_container(self, service_container: ServiceContainer | None) -> None:
        """Attach a service container after initialization."""
        self._service_container = service_container
        if service_container is not None and self._llm_entity_extractor is None:
            self._llm_entity_extractor = LLMEntityExtractor(service_container)
    
    def _init_extractors(self):
        """Initialize available extractors"""
        # Try PDF extractors
        try:
            import pdfplumber
            self._pdf_extractor = 'pdfplumber'
        except ImportError:
            try:
                import PyPDF2
                self._pdf_extractor = 'pypdf2'
            except ImportError:
                logger.info("No PDF extractor available (install pdfplumber or PyPDF2)")
        
        # Try DOCX extractor
        try:
            import docx
            self._docx_available = True
        except ImportError:
            logger.info("python-docx not available for DOCX extraction")
    
    def extract_content(self, file_path: Path, max_chars: int = 10000) -> Dict[str, Any]:
        """
        Extract content and metadata from a file.
        
        Returns:
            {
                'text': str,  # Extracted text content
                'doc_type': str,  # Detected document type
                'entities': dict,  # Extracted entities
                'quality_score': float,  # Filename quality score
                'suggested_name': str,  # Suggested filename if poor quality
            }
        """
        result = {
            'text': '',
            'doc_type': None,
            'entities': {},
            'quality_score': 0.0,
            'suggested_name': None,
            'extraction_method': None,
        }
        
        suffix = file_path.suffix.lower()
        
        # Extract text based on file type
        if suffix == '.pdf':
            result['text'] = self._extract_pdf(file_path, max_chars)
            result['extraction_method'] = 'pdf'
        elif suffix in ['.docx', '.doc']:
            result['text'] = self._extract_docx(file_path, max_chars)
            result['extraction_method'] = 'docx'
        elif suffix in ['.txt', '.md', '.log', '.csv', '.json']:
            result['text'] = self._extract_text(file_path, max_chars)
            result['extraction_method'] = 'text'
        elif suffix in ['.htm', '.html']:
            result['text'] = self._extract_html(file_path, max_chars)
            result['extraction_method'] = 'html'
        elif suffix in ['.xlsx', '.xls']:
            result['text'] = self._extract_excel(file_path, max_chars)
            result['extraction_method'] = 'excel'
        elif suffix in ['.pptx', '.ppt']:
            result['text'] = self._extract_pptx(file_path, max_chars)
            result['extraction_method'] = 'pptx'
        
        # Combine filename and content for analysis
        full_text = f"{file_path.stem}\n{result['text']}"
        
        # Detect document type from content
        result['doc_type'] = self._detect_doc_type(full_text)
        
        # Extract entities
        result['entities'] = self._extract_entities(full_text)
        
        # Assess filename quality
        result['quality_score'] = self._assess_filename_quality(file_path.name)
        
        # Generate suggested name if quality is poor
        if result['quality_score'] < 0.5:
            result['suggested_name'] = self._generate_suggested_name(
                file_path, result['doc_type'], result['entities']
            )
        
        return result
    
    def _extract_pdf(self, file_path: Path, max_chars: int) -> str:
        """Extract text from PDF with 10-second timeout protection"""
        if not self._pdf_extractor:
            return ''

        # Quick check: is this actually a PDF file?
        try:
            with open(file_path, 'rb') as f:
                header = f.read(8)
                if not header.startswith(b'%PDF-'):
                    logger.debug(f"File {file_path} does not have valid PDF header (starts with {header[:4]})")
                    return ''
        except Exception as e:
            logger.debug(f"Could not read header of {file_path}: {e}")
            return ''

        try:
            if self._pdf_extractor == 'pdfplumber':
                import pdfplumber
                from concurrent.futures import ThreadPoolExecutor, TimeoutError as FutureTimeoutError

                def extract_with_timeout():
                    """Extract PDF content with internal timeout handling"""
                    text = []
                    with pdfplumber.open(file_path) as pdf:
                        # Extract from first 2 pages only for quick classification
                        for i, page in enumerate(pdf.pages[:2]):
                            try:
                                page_text = page.extract_text() or ''
                                text.append(page_text)
                                # Stop if we have enough content
                                if sum(len(t) for t in text) > max_chars:
                                    break
                            except Exception as e:
                                logger.debug(f"Failed to extract page {i} from {file_path}: {e}")
                                continue
                    return ''.join(text)[:max_chars]

                # Run extraction with 10-second timeout
                try:
                    with ThreadPoolExecutor(max_workers=1) as executor:
                        future = executor.submit(extract_with_timeout)
                        return future.result(timeout=10)  # 10 second timeout
                except FutureTimeoutError:
                    logger.warning(f"PDF extraction timed out after 10s for {file_path}")
                    return ''
                except Exception as e:
                    logger.warning(f"PDF extraction failed for {file_path}: {e}")
                    return ''

        except Exception as e:
            logger.warning(f"PDF extraction setup failed for {file_path}: {e}")
            return ''

        return ''
    
    def _extract_docx(self, file_path: Path, max_chars: int) -> str:
        """Extract text from DOCX"""
        if not self._docx_available:
            return ''
        
        try:
            import docx
            doc = docx.Document(file_path)
            text = []
            for para in doc.paragraphs:
                text.append(para.text)
                if sum(len(t) for t in text) > max_chars:
                    break
            return '\n'.join(text)[:max_chars]
        except Exception as e:
            logger.debug(f"DOCX extraction failed for {file_path}: {e}")
        return ''
    
    def _extract_text(self, file_path: Path, max_chars: int) -> str:
        """Extract text from plain text files"""
        try:
            return file_path.read_text(encoding='utf-8', errors='ignore')[:max_chars]
        except Exception as e:
            logger.debug(f"Text extraction failed for {file_path}: {e}")
        return ''
    
    def _extract_html(self, file_path: Path, max_chars: int) -> str:
        """Extract text from HTML files"""
        try:
            content = file_path.read_text(encoding='utf-8', errors='ignore')
            # Simple HTML tag removal
            text = re.sub(r'<script[^>]*>.*?</script>', '', content, flags=re.DOTALL | re.IGNORECASE)
            text = re.sub(r'<style[^>]*>.*?</style>', '', text, flags=re.DOTALL | re.IGNORECASE)
            text = re.sub(r'<[^>]+>', ' ', text)
            text = re.sub(r'\s+', ' ', text)
            return text.strip()[:max_chars]
        except Exception as e:
            logger.debug(f"HTML extraction failed for {file_path}: {e}")
        return ''
    
    def _extract_excel(self, file_path: Path, max_chars: int) -> str:
        """Extract text from Excel files (xlsx, xls)"""
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            text = []
            total_chars = 0
            
            for sheet in wb.sheetnames[:3]:  # First 3 sheets
                ws = wb[sheet]
                text.append(f"=== Sheet: {sheet} ===")
                
                for row in ws.iter_rows(max_row=50):  # First 50 rows
                    row_text = []
                    for cell in row:
                        if cell.value is not None:
                            row_text.append(str(cell.value))
                    if row_text:
                        line = ' | '.join(row_text)
                        text.append(line)
                        total_chars += len(line)
                        if total_chars > max_chars:
                            break
                if total_chars > max_chars:
                    break
            
            wb.close()
            return '\n'.join(text)[:max_chars]
        except ImportError:
            logger.debug("openpyxl not available for Excel extraction")
        except Exception as e:
            logger.debug(f"Excel extraction failed for {file_path}: {e}")
        return ''
    
    def _extract_pptx(self, file_path: Path, max_chars: int) -> str:
        """Extract text from PowerPoint files (pptx, ppt)"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text = []
            total_chars = 0
            
            for i, slide in enumerate(prs.slides[:20]):  # First 20 slides
                slide_text = [f"=== Slide {i+1} ==="]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                        total_chars += len(shape.text)
                
                text.extend(slide_text)
                if total_chars > max_chars:
                    break
            
            return '\n'.join(text)[:max_chars]
        except ImportError:
            logger.debug("python-pptx not available for PowerPoint extraction")
        except Exception as e:
            logger.debug(f"PowerPoint extraction failed for {file_path}: {e}")
        return ''
    
    def _detect_doc_type(self, text: str) -> Optional[str]:
        """Detect document type from content"""
        text_lower = text.lower()
        
        scores = {}
        for doc_type, keywords in self.DOC_TYPE_KEYWORDS.items():
            score = sum(1 for kw in keywords if kw in text_lower)
            if score > 0:
                scores[doc_type] = score
        
        if scores:
            return max(scores, key=scores.get)
        return None
    
    def _extract_entities(self, text: str) -> Dict[str, List[str]]:
        """Extract entities from text"""
        entities = {}
        
        for entity_type, patterns in self.ENTITY_PATTERNS.items():
            matches = set()
            for pattern in patterns:
                try:
                    for match in re.finditer(pattern, text, re.IGNORECASE):
                        # Get all groups or full match
                        if match.groups():
                            for g in match.groups():
                                if g:
                                    matches.add(g.strip())
                        else:
                            matches.add(match.group(0).strip())
                except Exception:
                    pass
            
            if matches:
                entities[entity_type] = list(matches)[:5]  # Limit to 5 per type
        
        return entities

    def extract_entities_with_llm(self, text: str) -> Tuple[List[ExtractedEntity], float]:
        """Extract entities using an LLM if available."""
        if self._llm_entity_extractor is None:
            return [], 0.0
        return self._llm_entity_extractor.extract(text)
    
    def _assess_filename_quality(self, filename: str) -> float:
        """
        Assess the quality of a filename for organization purposes.
        
        Returns a score from 0 (poor) to 1 (good).
        """
        score = 0.5  # Start neutral
        stem = Path(filename).stem
        
        # Positive factors
        # Has a date
        if re.search(r'\d{4}[-_]\d{2}[-_]\d{2}|\d{2}[-_]\d{2}[-_]\d{4}', stem):
            score += 0.15
        
        # Has a case number
        if re.search(r'\d{2}[A-Z]{2}\d{5,7}|\d{2,4}[-]?[A-Z]{0,3}[-]?\d{3,7}', stem):
            score += 0.15
        
        # Has meaningful words (not just numbers/codes)
        words = re.findall(r'[a-zA-Z]{3,}', stem)
        if len(words) >= 2:
            score += 0.1
        
        # Has document type keyword
        stem_lower = stem.lower()
        doc_types = ['motion', 'order', 'complaint', 'brief', 'affidavit', 'petition', 
                     'notice', 'report', 'letter', 'invoice', 'receipt', 'contract']
        if any(dt in stem_lower for dt in doc_types):
            score += 0.1
        
        # Negative factors
        # Very short name
        if len(stem) < 5:
            score -= 0.2
        
        # Mostly numbers or random characters
        if len(re.findall(r'[a-zA-Z]', stem)) < len(stem) * 0.3:
            score -= 0.15
        
        # Has version suffix like (2), (3), copy
        if re.search(r'\(\d+\)|copy|копия', stem, re.IGNORECASE):
            score -= 0.1
        
        # Very generic name
        generic_names = ['document', 'file', 'untitled', 'new', 'scan', 'img', 'doc']
        if stem_lower in generic_names or any(stem_lower.startswith(g) for g in generic_names):
            score -= 0.2
        
        # Random hash-like names
        if re.match(r'^[a-f0-9]{8,}$', stem, re.IGNORECASE):
            score -= 0.3
        
        return max(0.0, min(1.0, score))
    
    def _generate_suggested_name(
        self, 
        file_path: Path, 
        doc_type: Optional[str], 
        entities: Dict[str, List[str]]
    ) -> Optional[str]:
        """Generate a suggested filename based on extracted information"""
        parts = []
        
        # Add date
        if 'date' in entities and entities['date']:
            date_str = entities['date'][0]
            # Try to normalize date format
            try:
                # Try different formats
                for fmt in ['%m/%d/%Y', '%Y-%m-%d', '%B %d, %Y', '%m/%d/%y']:
                    try:
                        dt = datetime.strptime(date_str, fmt)
                        parts.append(dt.strftime('%Y-%m-%d'))
                        break
                    except ValueError:
                        continue
            except:
                pass
        
        # Add case number
        if 'case_number' in entities and entities['case_number']:
            parts.append(entities['case_number'][0].replace(' ', ''))
        
        # Add document type
        if doc_type:
            parts.append(doc_type.capitalize())
        
        # Add court if available
        if 'court' in entities and entities['court']:
            court = entities['court'][0]
            # Abbreviate court name
            court_abbrev = ''.join(word[0] for word in court.split() if word[0].isupper())
            if court_abbrev:
                parts.append(court_abbrev)
        
        # Add parties if available
        if 'parties' in entities and entities['parties']:
            party = entities['parties'][0][:20]  # Limit length
            party = re.sub(r'[^\w\s]', '', party).strip()
            if party:
                parts.append(party.replace(' ', '_'))
        
        if len(parts) >= 2:
            return '_'.join(parts) + file_path.suffix
        
        return None
    
    def quick_classify(self, file_path: Path) -> Dict[str, Any]:
        """
        Quick classification using only filename and first few KB of content.
        Faster than full extraction.
        """
        result = {
            'doc_type': None,
            'entities': {},
            'quality_score': 0.0,
        }
        
        # Analyze filename
        filename = file_path.name
        stem = file_path.stem
        
        # Quality score
        result['quality_score'] = self._assess_filename_quality(filename)
        
        # Extract entities from filename
        result['entities'] = self._extract_entities(stem)
        
        # Detect doc type from filename
        result['doc_type'] = self._detect_doc_type(stem)
        
        # If filename isn't informative, try to read first part of file
        if result['quality_score'] < 0.5 or not result['doc_type']:
            try:
                suffix = file_path.suffix.lower()
                sample = ''

                try:
                    if suffix in ['.txt', '.md', '.log', '.csv']:
                        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                            sample = f.read(2000)
                    elif suffix == '.pdf' and self._pdf_extractor:
                        sample = self._extract_pdf(file_path, 2000)
                    elif suffix in ['.docx'] and self._docx_available:
                        sample = self._extract_docx(file_path, 2000)
                except Exception as e:
                    logger.debug(f"Failed to sample content from {file_path}: {e}")
                    sample = ''

                if sample:
                    if not result['doc_type']:
                        result['doc_type'] = self._detect_doc_type(sample)

            except Exception as e:
                logger.debug(f"Content sampling failed for {file_path}: {e}")
                # Continue without content sampling

        return result


def get_extractor(service_container: ServiceContainer | None = None) -> ContentExtractor:
    """Get singleton content extractor instance."""
    if not hasattr(get_extractor, '_instance'):
        get_extractor._instance = ContentExtractor(service_container=service_container)
        return get_extractor._instance

    if service_container is not None:
        instance = get_extractor._instance
        instance.set_service_container(service_container)

    return get_extractor._instance
