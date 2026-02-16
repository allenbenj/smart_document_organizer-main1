"""
Unified Document Processor Agent - Production Legal AI
======================================================

A comprehensive document processing agent that consolidates the best features from
existing implementations while maintaining collective intelligence integration.

This agent processes legal documents across all formats with intelligent strategy
selection and integrates with the shared memory system for collective intelligence.

Key Features:
- Multi-format processing: PDF, DOCX, TXT, HTML, RTF, MD, XLSX, CSV, PPTX, Images
- Intelligent processing strategies based on document type
- OCR fallback for scanned documents and images
- Multimodal capabilities (video depositions, audio transcription)
- Structured data extraction (Excel, CSV to database tables)
- Legal document classification and metadata extraction
- Integration with shared memory for collective intelligence
- Production-grade dependency management with graceful fallbacks
- Async processing with comprehensive error handling

Consolidates functionality from:
- document_processor_agent.py (dependency management)
- document_processor_agent_v2.py (multimodal capabilities)  
- document_processor_full.py (comprehensive format support)
"""

import asyncio
import hashlib
import io
import json
import logging
import mimetypes
import os
import re
import tempfile
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional, Union

# Core imports
# These are placeholder imports for the example. Replace with your actual project structure.
from ..base import BaseAgent, AgentResult, AgentStatus, TaskPriority
from ..base.agent_mixins import DocumentProcessingMixin, MemoryEnabledMixin
from ...core.container.service_container_impl import ProductionServiceContainer
try:
    from ...memory import MemoryType  # type: ignore
except Exception:  # Fallback stub when memory module is unavailable
    class MemoryType:  # type: ignore
        DOCUMENT = "document"

# Optional dependencies with graceful degradation
try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    fitz = None
    PYMUPDF_AVAILABLE = False

try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DocxDocument = None
    DOCX_AVAILABLE = False

try:
    from bs4 import BeautifulSoup
    BS4_AVAILABLE = True
except ImportError:
    BeautifulSoup = None
    BS4_AVAILABLE = False

try:
    import pytesseract
    from PIL import Image
    OCR_AVAILABLE = True
except ImportError:
    pytesseract = None
    Image = None
    OCR_AVAILABLE = False

try:
    import pandas as pd
    import openpyxl
    EXCEL_AVAILABLE = True
except ImportError:
    pd = None
    openpyxl = None
    EXCEL_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    Presentation = None
    PPTX_AVAILABLE = False

try:
    import markdown
    MARKDOWN_AVAILABLE = True
except ImportError:
    markdown = None
    MARKDOWN_AVAILABLE = False

try:
    from striprtf.striprtf import rtf_to_text
    RTF_AVAILABLE = True
except ImportError:
    rtf_to_text = None
    RTF_AVAILABLE = False

# Multimodal processing dependencies
try:
    from moviepy.editor import VideoFileClip
    import whisper
    MULTIMODAL_AVAILABLE = True
except ImportError:
    VideoFileClip = None
    whisper = None
    MULTIMODAL_AVAILABLE = False

logger = logging.getLogger(__name__)


class ProcessingStrategy:
    """Document processing strategy enumeration"""
    FULL_PROCESSING = "full_processing"  # Complete text extraction and analysis
    STRUCTURED_DATA = "structured_data"  # Extract to database tables (Excel, CSV)
    REFERENCE_ONLY = "reference_only"    # Metadata + basic text only (PPTX, Images)
    MULTIMODAL = "multimodal"           # Video/audio with transcription
    OCR_ONLY = "ocr_only"               # For images or scanned PDFs
    METADATA_ONLY = "metadata_only"     # Only extract file system and basic metadata


class DocumentFormat:
    """Supported document format enumeration"""
    # Text documents (full processing)
    PDF = "pdf"
    DOCX = "docx"
    TXT = "txt"
    MD = "markdown"
    HTML = "html"
    RTF = "rtf"
    
    # Structured data (database extraction)
    XLSX = "xlsx"
    CSV = "csv"
    
    # Presentation (reference only)
    PPTX = "pptx"
    
    # Images (OCR if needed)
    PNG = "png"
    JPG = "jpg"
    JPEG = "jpeg"
    TIFF = "tiff"
    BMP = "bmp"
    
    # Multimodal
    MP4 = "mp4"
    MOV = "mov"
    AVI = "avi"
    MP3 = "mp3"
    WAV = "wav"

    # Unsupported but recognized
    DOC = "doc"
    XLS = "xls"
    PPT = "ppt"
    
    UNKNOWN = "unknown"


@dataclass
class DocumentProcessingConfig:
    """Configuration for document processing"""
    max_file_size_mb: int = 100
    enable_ocr: bool = True
    enable_metadata_extraction: bool = True
    enable_chunking: bool = True
    chunk_size: int = 1000
    chunk_overlap: int = 200
    supported_formats: List[str] = field(default_factory=lambda: [
        '.pdf', '.docx', '.doc', '.txt', '.html', '.htm', '.rtf', '.md',
        '.xlsx', '.xls', '.csv', '.pptx', '.ppt',
        '.png', '.jpg', '.jpeg', '.tiff', '.bmp',
        '.mp4', '.mov', '.avi', '.mp3', '.wav'
    ])
    ocr_language: str = 'eng'
    extract_images_from_pdf: bool = False # Feature flag for image extraction
    enable_multimodal: bool = True
    enable_structured_data: bool = True
    whisper_model: str = 'base'  # Whisper model for audio transcription
    max_audio_duration_minutes: int = 60
    legal_document_detection: bool = True


@dataclass
class ProcessedDocument:
    """Represents a processed document with enhanced multimodal support"""
    document_id: str
    file_path: str
    filename: str
    file_size: int
    mime_type: str
    content: str
    metadata: Dict[str, Any]
    chunks: List[Dict[str, Any]]
    images: List[Dict[str, Any]]
    processing_time: float
    processing_method: str
    confidence: float
    document_format: str = DocumentFormat.UNKNOWN
    processing_strategy: str = ProcessingStrategy.FULL_PROCESSING
    structured_data: Optional[List[Dict[str, Any]]] = None  # For Excel/CSV data
    transcription: Optional[str] = None  # For audio/video content
    legal_indicators: List[str] = field(default_factory=list)  # Legal document markers
    multimodal_components: Dict[str, Any] = field(default_factory=dict)  # Additional components
    
    def to_dict(self) -> Dict[str, Any]:
        """Convert to dictionary for storage"""
        return {
            'document_id': self.document_id,
            'file_path': self.file_path,
            'filename': self.filename,
            'file_size': self.file_size,
            'mime_type': self.mime_type,
            'content': self.content,
            'metadata': self.metadata,
            'chunks': self.chunks,
            'images': self.images,
            'processing_time': self.processing_time,
            'processing_method': self.processing_method,
            'confidence': self.confidence,
            'document_format': self.document_format,
            'processing_strategy': self.processing_strategy,
            'structured_data': self.structured_data,
            'transcription': self.transcription,
            'legal_indicators': self.legal_indicators,
            'multimodal_components': self.multimodal_components
        }


class DocumentProcessor(BaseAgent, DocumentProcessingMixin, MemoryEnabledMixin):
    """
    Production Document Processor with collective intelligence.
    
    Processes various document formats and contributes to the shared
    knowledge base that enables collective intelligence across agents.
    """
    
    def __init__(
        self,
        services: ProductionServiceContainer,
        config: Optional[DocumentProcessingConfig] = None
    ):
        # Initialize base agent
        super().__init__(
            services=services,
            agent_name="DocumentProcessor",
            agent_type="document_processing",
            timeout_seconds=900.0  # Document processing can take time
        )
        
        # Initialize mixins
        DocumentProcessingMixin.__init__(self)
        MemoryEnabledMixin.__init__(self)
        
        # Configuration
        self.config = config or DocumentProcessingConfig()
        
        # Processing statistics
        self.stats = {
            'documents_processed': 0,
            'total_size_processed_mb': 0.0,
            'formats_processed': {},
            'avg_processing_time': 0.0,
            'total_processing_time': 0.0,
            'ocr_used_count': 0,
            'errors_encountered': 0
        }
        
        # Document format handlers
        self.format_handlers = {
            '.pdf': self._process_pdf,
            '.docx': self._process_docx,
            '.doc': self._process_docx,  # Will raise NotImplementedError
            '.txt': self._process_text,
            '.html': self._process_html,
            '.htm': self._process_html,
            '.rtf': self._process_rtf,
            '.md': self._process_markdown,
            '.xlsx': self._process_excel,
            '.xls': self._process_excel,  # Will raise NotImplementedError
            '.csv': self._process_csv,
            '.pptx': self._process_powerpoint,
            '.ppt': self._process_powerpoint, # Will raise NotImplementedError
            '.png': self._process_image,
            '.jpg': self._process_image,
            '.jpeg': self._process_image,
            '.tiff': self._process_image,
            '.bmp': self._process_image,
            '.mp4': self._process_video,
            '.mov': self._process_video,
            '.avi': self._process_video,
            '.mp3': self._process_audio,
            '.wav': self._process_audio
        }
        
        # Legal document indicators
        self.legal_indicators = [
            'motion', 'complaint', 'affidavit', 'deposition', 'subpoena',
            'warrant', 'indictment', 'plea', 'sentence', 'judgment',
            'order', 'ruling', 'opinion', 'brief', 'memorandum',
            'constitution', 'amendment', 'statute', 'regulation',
            'violation', 'misconduct', 'evidence', 'testimony',
            'plaintiff', 'defendant', 'appellant', 'appellee',
            'contract', 'agreement', 'lease', 'merger', 'acquisition'
        ]
        
        logger.info(f"DocumentProcessor initialized with config: {self.config}")
        self._log_available_dependencies()
    
    def _log_available_dependencies(self):
        """Log available optional dependencies."""
        dependencies = {
            'PyMuPDF (PDF)': PYMUPDF_AVAILABLE,
            'python-docx (DOCX)': DOCX_AVAILABLE,
            'BeautifulSoup (HTML)': BS4_AVAILABLE,
            'pytesseract (OCR)': OCR_AVAILABLE,
            'pandas (Excel/CSV)': EXCEL_AVAILABLE,
            'python-pptx (PowerPoint)': PPTX_AVAILABLE,
            'markdown (Markdown)': MARKDOWN_AVAILABLE,
            'striprtf (RTF)': RTF_AVAILABLE,
            'whisper + moviepy (Multimodal)': MULTIMODAL_AVAILABLE
        }
        available = [name for name, avail in dependencies.items() if avail]
        unavailable = [name for name, avail in dependencies.items() if not avail]
        logger.info(f"Available dependencies: {', '.join(available) if available else 'None'}")
        if unavailable:
            logger.warning(f"Unavailable dependencies: {', '.join(unavailable)}. Some file types may not be processed.")

    async def _process_task(self, task_data: Any, metadata: Dict[str, Any]) -> Dict[str, Any]:
        """Process a document processing task."""
        file_path, content, document_id = None, None, None
        try:
            # Step 1: Resolve input data
            if isinstance(task_data, str):
                # Ambiguous case: could be a path or content.
                # Heuristic: treat as a path if it exists, otherwise content.
                p = Path(task_data)
                if p.exists() and p.is_file():
                    file_path = p
                    document_id = metadata.get('document_id', f"doc_{file_path.stem}_{file_path.stat().st_size}")
                else:
                    content = task_data
                    content_hash = hashlib.sha256(content.encode('utf-8')).hexdigest()
                    document_id = metadata.get('document_id', f"content_{content_hash[:16]}")
            elif isinstance(task_data, dict):
                file_path = Path(task_data['file_path']) if 'file_path' in task_data else None
                content = task_data.get('content')
                document_id = task_data.get('document_id', metadata.get('document_id'))
            else:
                raise ValueError(f"Unsupported task_data type: {type(task_data)}")

            if not file_path and not content:
                raise ValueError("No file path or content provided for document processing.")
            
            logger.info(f"Starting document processing for {document_id}")
            start_time = datetime.now()

            # Step 2: Process the document
            if file_path:
                processed_doc = await self._process_file(file_path, document_id)
            else:
                processed_doc = await self._process_content(content, document_id)
            
            # Step 3: Integrate with collective intelligence
            if self._is_memory_available():
                similar_docs = await self._find_similar_documents(processed_doc)
                if similar_docs:
                    logger.info(f"Found {len(similar_docs)} similar documents. Enhancing...")
                    processed_doc = self._enhance_with_collective_intelligence(processed_doc, similar_docs)
                await self._store_processed_document(processed_doc)
            else:
                similar_docs = []
            
            # Step 4: Update stats and finalize
            self._update_statistics(processed_doc)
            processing_time = (datetime.now() - start_time).total_seconds()
            processed_doc.processing_time = processing_time # Update with total time

            logger.info(f"Document processing completed for {document_id} in {processing_time:.2f}s")
            
            return {
                'success': True,
                'processed_document': processed_doc.to_dict(),
                'collective_intelligence': {
                    'similar_documents_found': len(similar_docs),
                    'knowledge_enhanced': len(similar_docs) > 0
                }
            }
            
        except Exception as e:
            self.stats['errors_encountered'] += 1
            doc_id = document_id or metadata.get('document_id', 'unknown')
            logger.error(f"Document processing failed for {doc_id}: {e}", exc_info=True)
            raise

    # --- Core Processing Logic ---

    async def _process_file(self, file_path: Path, document_id: str) -> ProcessedDocument:
        """Process a file from the filesystem."""
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        file_size = file_path.stat().st_size
        max_size = self.config.max_file_size_mb * 1024 * 1024
        if file_size > max_size:
            raise ValueError(f"File too large: {file_size / (1024*1024):.1f}MB > {self.config.max_file_size_mb}MB")
        
        mime_type, _ = mimetypes.guess_type(str(file_path))
        file_extension = file_path.suffix.lower()
        
        if file_extension not in self.config.supported_formats:
            raise ValueError(f"Unsupported file format: {file_extension}")
        
        logger.info(f"Processing file: {file_path} (ext: {file_extension}, size: {file_size} bytes)")
        
        handler = self.format_handlers.get(file_extension)
        if not handler:
            raise ValueError(f"No handler available for format: {file_extension}")
        
        processing_strategy = self._determine_processing_strategy(file_extension)
        document_format = self._get_document_format(file_extension)
        
        start_time = datetime.now()
        content, metadata, structured, transcription, processing_method = "", {}, None, None, "unknown"
        try:
            content, metadata, structured, transcription, processing_method = await handler(file_path)
            confidence = self._calculate_processing_confidence(processing_method, content, file_extension)
        except Exception as e:
            logger.warning(f"Handler for {file_extension} failed: {e}. Attempting fallback.")
            content, metadata, processing_method = await self._fallback_processing(file_path)
            confidence = 0.3
        
        processing_time = (datetime.now() - start_time).total_seconds()
        
        legal_indicators = self._detect_legal_content(content or transcription or "", file_path.name)
        chunks = self._create_chunks(content) if self.config.enable_chunking and content else []
        
        return ProcessedDocument(
            document_id=document_id,
            file_path=str(file_path),
            filename=file_path.name,
            file_size=file_size,
            mime_type=mime_type or 'application/octet-stream',
            content=content,
            metadata=metadata,
            chunks=chunks,
            images=[],  # Image extraction is a separate, complex feature
            processing_time=processing_time,
            processing_method=processing_method,
            confidence=confidence,
            document_format=document_format,
            processing_strategy=processing_strategy,
            legal_indicators=legal_indicators,
            structured_data=structured,
            transcription=transcription
        )

    async def _process_content(self, content: str, document_id: str) -> ProcessedDocument:
        """Process raw text content."""
        start_time = datetime.now()
        chunks = self._create_chunks(content) if self.config.enable_chunking else []
        legal_indicators = self._detect_legal_content(content, f"{document_id}.txt")
        processing_time = (datetime.now() - start_time).total_seconds()

        return ProcessedDocument(
            document_id=document_id,
            file_path="",
            filename=f"{document_id}.txt",
            file_size=len(content.encode('utf-8')),
            mime_type="text/plain",
            content=content,
            metadata={'source': 'direct_content'},
            chunks=chunks,
            images=[],
            processing_time=processing_time,
            processing_method="direct_content",
            confidence=1.0,
            document_format=DocumentFormat.TXT,
            processing_strategy=ProcessingStrategy.FULL_PROCESSING,
            legal_indicators=legal_indicators
        )

    # --- File Type Handlers ---

    async def _process_pdf(self, file_path: Path) -> tuple:
        if not PYMUPDF_AVAILABLE:
            raise RuntimeError("PyMuPDF (fitz) is not installed, cannot process PDF files.")
        
        content, metadata, processing_method = "", {}, "pymupdf"
        doc = fitz.open(str(file_path))
        metadata = doc.metadata
        metadata['page_count'] = doc.page_count
        metadata['encrypted'] = doc.is_encrypted

        for page_num, page in enumerate(doc):
            page_text = page.get_text()
            if page_text.strip():
                content += f"\n--- Page {page_num + 1} ---\n{page_text}"
            elif self.config.enable_ocr and OCR_AVAILABLE:
                try:
                    pix = page.get_pixmap()
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    ocr_text = pytesseract.image_to_string(img, lang=self.config.ocr_language)
                    if ocr_text.strip():
                        content += f"\n--- Page {page_num + 1} (OCR) ---\n{ocr_text}"
                        processing_method = "pymupdf_with_ocr"
                        self.stats['ocr_used_count'] += 1
                except Exception as e:
                    logger.warning(f"OCR failed for page {page_num + 1} in {file_path.name}: {e}")
        doc.close()
        return content.strip(), metadata, None, None, processing_method

    async def _process_docx(self, file_path: Path) -> tuple:
        if file_path.suffix.lower() == '.doc':
            raise NotImplementedError(".doc files are not supported. Please convert to .docx first.")
        if not DOCX_AVAILABLE:
            raise RuntimeError("python-docx is not installed, cannot process DOCX files.")
        
        doc = DocxDocument(str(file_path))
        props = doc.core_properties
        metadata = {
            'author': props.author, 'created': props.created, 'modified': props.modified,
            'title': props.title, 'subject': props.subject
        }
        full_text = [p.text for p in doc.paragraphs]
        for table in doc.tables:
            for row in table.rows:
                full_text.append(" | ".join(cell.text for cell in row.cells))
        return "\n".join(full_text), metadata, None, None, "python-docx"

    async def _process_text(self, file_path: Path) -> tuple:
        content, encoding_used = "", None
        for encoding in ['utf-8', 'latin-1', 'cp1252']:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    content = f.read()
                encoding_used = encoding
                break
            except UnicodeDecodeError:
                continue
        if not encoding_used:
            raise ValueError("Could not decode text file with any supported encoding.")
        metadata = {'encoding': encoding_used}
        return content, metadata, None, None, "text_read"

    async def _process_html(self, file_path: Path) -> tuple:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            html_content = f.read()
        if BS4_AVAILABLE:
            soup = BeautifulSoup(html_content, 'html.parser')
            for script in soup(["script", "style"]):
                script.extract()
            content = soup.get_text(separator='\n', strip=True)
            title = soup.title.string if soup.title else ""
            return content, {'title': title}, None, None, "beautifulsoup"
        else:
            content = re.sub(r'<[^>]+>', '', html_content)
            return content, {}, None, None, "regex_html_strip"

    async def _process_rtf(self, file_path: Path) -> tuple:
        with open(file_path, 'r', errors='ignore') as f:
            rtf_content = f.read()
        if RTF_AVAILABLE:
            content = rtf_to_text(rtf_content)
            return content, {}, None, None, "striprtf"
        else:
            # Basic fallback if striprtf is not available
            content = re.sub(r'{\\*\\.+?}|\\(.+?)\s|[{}]', '', rtf_content)
            return content.strip(), {}, None, None, "regex_rtf_strip"

    async def _process_markdown(self, file_path: Path) -> tuple:
        with open(file_path, 'r', encoding='utf-8') as f:
            md_content = f.read()
        if MARKDOWN_AVAILABLE and BS4_AVAILABLE:
            html = markdown.markdown(md_content)
            soup = BeautifulSoup(html, 'html.parser')
            content = soup.get_text()
            return content, {}, None, None, "markdown_to_text"
        else:
            # Fallback for when markdown or bs4 is not available
            content = re.sub(r'#+\s*|\*\*|\*|`|\[.+?\]\(.+?\)', '', md_content)
            return content, {}, None, None, "regex_markdown_strip"

    async def _process_excel(self, file_path: Path) -> tuple:
        if file_path.suffix.lower() == '.xls':
            raise NotImplementedError(".xls files are not supported. Please convert to .xlsx first.")
        if not EXCEL_AVAILABLE:
            raise RuntimeError("pandas/openpyxl not installed, cannot process Excel files.")
        
        xls = pd.ExcelFile(str(file_path))
        structured_data, text_content = [], []
        for sheet_name in xls.sheet_names:
            df = pd.read_excel(xls, sheet_name=sheet_name)
            df = df.dropna(how='all').fillna('') # Clean up empty rows and NaN
            if df.empty: continue
            
            # Convert to list of dicts for structured data
            sheet_data = df.to_dict(orient='records')
            structured_data.extend(sheet_data)
            
            # Create a textual representation for content field
            text_content.append(f"--- Sheet: {sheet_name} ---\n")
            text_content.append(df.to_string(index=False))
            
        content = "\n".join(text_content)
        metadata = {'sheets': xls.sheet_names, 'total_rows': len(structured_data)}
        return content, metadata, structured_data, None, "pandas_excel"

    async def _process_csv(self, file_path: Path) -> tuple:
        if not EXCEL_AVAILABLE: # pandas is used for CSV too
            raise RuntimeError("pandas not installed, cannot process CSV files.")
            
        df = pd.read_csv(str(file_path)).fillna('')
        structured_data = df.to_dict(orient='records')
        content = df.to_string(index=False)
        metadata = {'columns': list(df.columns), 'total_rows': len(structured_data)}
        return content, metadata, structured_data, None, "pandas_csv"

    async def _process_powerpoint(self, file_path: Path) -> tuple:
        if file_path.suffix.lower() == '.ppt':
            raise NotImplementedError(".ppt files are not supported. Please convert to .pptx first.")
        if not PPTX_AVAILABLE:
            raise RuntimeError("python-pptx not installed, cannot process PowerPoint files.")
            
        prs = Presentation(str(file_path))
        text_runs = []
        for i, slide in enumerate(prs.slides):
            text_runs.append(f"\n--- Slide {i+1} ---\n")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        content = "\n".join(text_runs)
        metadata = {'slide_count': len(prs.slides)}
        return content, metadata, None, None, "python-pptx"

    async def _process_image(self, file_path: Path) -> tuple:
        if not OCR_AVAILABLE:
            raise RuntimeError("pytesseract/Pillow not installed, cannot process images.")
        
        img = Image.open(str(file_path))
        content = pytesseract.image_to_string(img, lang=self.config.ocr_language)
        self.stats['ocr_used_count'] += 1
        metadata = {'image_size': img.size, 'image_mode': img.mode}
        return content, metadata, None, None, "tesseract_ocr"

    async def _process_video(self, file_path: Path) -> tuple:
        if not MULTIMODAL_AVAILABLE:
            raise RuntimeError("moviepy not installed, cannot process video files.")

        with tempfile.TemporaryDirectory() as tmpdir:
            audio_path = Path(tmpdir) / f"{file_path.stem}.wav"
            try:
                video = VideoFileClip(str(file_path))
                if video.duration > self.config.max_audio_duration_minutes * 60:
                     raise ValueError("Video duration exceeds maximum limit.")
                video.audio.write_audiofile(str(audio_path), codec='pcm_s16le')
                video.close()
                return await self._process_audio(audio_path)
            except Exception as e:
                logger.error(f"Failed to extract audio from video {file_path.name}: {e}")
                return "", {"error": str(e)}, None, None, "video_audio_extraction_failed"

    async def _process_audio(self, file_path: Path) -> tuple:
        if not MULTIMODAL_AVAILABLE:
            raise RuntimeError("whisper not installed, cannot process audio files.")

        model = whisper.load_model(self.config.whisper_model)
        result = model.transcribe(str(file_path))
        transcription = result['text']
        metadata = {'language': result['language'], 'model': self.config.whisper_model}
        return "", metadata, None, transcription, "whisper_transcription"

    async def _fallback_processing(self, file_path: Path) -> tuple:
        """Fallback processing for when specific handlers fail."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            metadata = {
                'fallback_method': 'text_read_ignore_errors',
                'warning': 'Processed with fallback method; content may be incomplete or garbled.'
            }
            return content, metadata, "fallback_text"
        except Exception as e:
            logger.error(f"Final fallback processing failed for {file_path}: {e}")
            return "", {'error': 'Could not process file with any method'}, "fallback_failed"

    # --- Helper & Integration Methods ---
    
    def _determine_processing_strategy(self, file_extension: str) -> str:
        """Determines the processing strategy based on file extension."""
        if file_extension in ['.xlsx', '.xls', '.csv']:
            return ProcessingStrategy.STRUCTURED_DATA
        elif file_extension in ['.mp4', '.mov', '.avi', '.mp3', '.wav']:
            return ProcessingStrategy.MULTIMODAL
        elif file_extension in ['.png', '.jpg', '.jpeg', '.tiff', '.bmp']:
            return ProcessingStrategy.OCR_ONLY
        elif file_extension in ['.pptx', '.ppt']:
            return ProcessingStrategy.REFERENCE_ONLY
        else:
            return ProcessingStrategy.FULL_PROCESSING

    def _get_document_format(self, file_extension: str) -> str:
        """Gets the document format enum from the file extension."""
        return getattr(DocumentFormat, file_extension.lstrip('.').upper(), DocumentFormat.UNKNOWN)

    def _calculate_processing_confidence(self, method: str, content: str, ext: str) -> float:
        """Calculates a confidence score for the processing."""
        if not content: return 0.1
        if "failed" in method: return 0.0
        if "fallback" in method: return 0.3
        if "ocr" in method: return 0.75
        if "regex" in method: return 0.6
        return 0.95 # High confidence for direct library-based extraction

    def _create_chunks(self, text: str) -> List[Dict[str, Any]]:
        """Splits text into chunks with overlap."""
        chunks = []
        start = 0
        doc_id_base = hashlib.sha256(text.encode('utf-8')).hexdigest()[:8]
        while start < len(text):
            end = start + self.config.chunk_size
            chunk_text = text[start:end]
            chunks.append({
                'chunk_id': f'chk_{doc_id_base}_{start}',
                'text': chunk_text,
                'start_index': start
            })
            start += self.config.chunk_size - self.config.chunk_overlap
        return chunks

    def _detect_legal_content(self, content: str, filename: str) -> List[str]:
        """Detects keywords indicating a legal document."""
        if not self.config.legal_document_detection: return []
        
        found_indicators = set()
        text_to_search = (content + " " + filename).lower()
        for indicator in self.legal_indicators:
            if indicator in text_to_search:
                found_indicators.add(indicator)
        return sorted(list(found_indicators))

    def _update_statistics(self, doc: ProcessedDocument):
        """Updates the agent's processing statistics."""
        self.stats['documents_processed'] += 1
        size_mb = doc.file_size / (1024 * 1024)
        self.stats['total_size_processed_mb'] += size_mb
        
        self.stats['formats_processed'][doc.document_format] = self.stats['formats_processed'].get(doc.document_format, 0) + 1
        
        total_time = self.stats['total_processing_time'] + doc.processing_time
        self.stats['total_processing_time'] = total_time
        self.stats['avg_processing_time'] = total_time / self.stats['documents_processed']

    async def _find_similar_documents(self, doc: ProcessedDocument) -> List[Any]:
        """Finds similar documents from shared memory."""
        query = " ".join(doc.content.split()[:50]) # Use first 50 words as query
        if not query: return []
        
        try:
            return await self.search_memory(
                query=query,
                memory_types=[MemoryType.DOCUMENT],
                limit=3,
                min_similarity=0.7
            )
        except Exception as e:
            logger.warning(f"Failed to search memory for similar documents: {e}")
            return []

    def _enhance_with_collective_intelligence(self, doc: ProcessedDocument, similar_docs: List[Any]) -> ProcessedDocument:
        """Enhances a document with data from similar documents."""
        collective_metadata = {'similar_document_ids': [d.id for d in similar_docs]}
        # Example: Inherit tags or classification from highly similar documents
        # This part would contain more complex logic based on your memory structure
        doc.metadata['collective_intelligence'] = collective_metadata
        return doc

    async def _store_processed_document(self, doc: ProcessedDocument):
        """Stores the processed document in the shared memory system."""
        try:
            await self.add_to_memory(
                memory_type=MemoryType.DOCUMENT,
                content=doc.content,
                metadata=doc.to_dict(), # Store the full object as metadata
                document_id=doc.document_id
            )
            logger.info(f"Stored document {doc.document_id} in collective memory.")
        except Exception as e:
            logger.error(f"Failed to store document {doc.document_id} in memory: {e}")


async def create_document_processor(services: ProductionServiceContainer, config: Optional[DocumentProcessingConfig] = None) -> DocumentProcessor:
    """Factory function to create a DocumentProcessor for the ProductionAgentManager.

    Accepts a service container and optional config. Returns an initialized
    DocumentProcessor instance ready to handle tasks.
    """
    return DocumentProcessor(services=services, config=config)
